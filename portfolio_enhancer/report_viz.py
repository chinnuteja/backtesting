# portfolio_enhancer/report_viz.py
from __future__ import annotations

import os
from typing import Dict, Optional, Tuple
import pandas as pd
import numpy as np
import matplotlib.pyplot as plt

# Optional PowerPoint export (graceful if missing)
try:
    from pptx import Presentation
    from pptx.util import Inches, Pt
    _PPTX_AVAILABLE = True
except Exception:
    _PPTX_AVAILABLE = False


# ==========
# Helpers
# ==========
def _ensure_dt_index(df: pd.DataFrame, col: str) -> pd.DataFrame:
    out = df.copy()
    if col in out.columns:
        out["__dt__"] = pd.to_datetime(out[col])
        out = out.sort_values("__dt__").set_index("__dt__")
        out = out.drop(columns=[col])
    if not isinstance(out.index, pd.DatetimeIndex):
        raise ValueError("log_df must contain a datetime index or a parseable period_end/date column")
    out.index = out.index.tz_localize(None)
    return out

def _year_slice(df: pd.DataFrame, year: int) -> pd.DataFrame:
    start, end = pd.Timestamp(f"{year}-01-01"), pd.Timestamp(f"{year}-12-31")
    return df.loc[(df.index >= start) & (df.index <= end)]

def _shade_event_band(ax, start: str, end: str, label: str):
    s, e = pd.Timestamp(start), pd.Timestamp(end)
    ax.axvspan(s, e, alpha=0.08, linewidth=0)
    # place label top-left of band
    y_max = ax.get_ylim()[1]
    ax.text(s, y_max, f" {label}", va="top", ha="left", fontsize=9)

def _add_q3_2023_band(ax):
    _shade_event_band(ax, "2023-07-01", "2023-09-30", "Q3'23: Rising US yields / Strong USD / Oil ↑")

def _kpi_bands(ax):
    # visual thresholds to make KPI scale intuitive
    ax.axhline(+0.50, linewidth=1, alpha=0.3)
    ax.axhline(0.00,  linewidth=1, alpha=0.6)
    ax.axhline(-0.50, linewidth=1, alpha=0.3)

def _pct_fmt(ax):
    ax.yaxis.set_major_formatter(lambda y, pos: f"{100*y:.0f}%")

def _add_footer(fig, txt: str):
    fig.text(0.01, 0.01, txt, fontsize=8, alpha=0.6)

def _summary_stats(real_series: pd.Series) -> Tuple[float, float, float, float]:
    real = real_series.fillna(0.0)
    total = float((1 + real).prod() - 1.0)
    vol   = float(real.std(ddof=1) * np.sqrt(12.0)) if len(real) > 1 else 0.0
    sharpe = (total / vol) if vol > 1e-12 else 0.0
    maxdd = 0.0
    cum = (1 + real).cumprod()
    peak = cum.expanding().max()
    dd = cum / peak - 1.0
    maxdd = float(dd.min()) if len(dd) else 0.0
    return total, vol, sharpe, maxdd


# ==========
# Core builder
# ==========
def build_visual_report(
    log_df: pd.DataFrame,
    daily_returns: Optional[pd.DataFrame] = None,
    title: str = "2023 Portfolio Forecast Report",
    outfile_prefix: str = "Portfolio_Report_2023",
    focus_year: int = 2023,
):
    """
    Inputs
    ------
    log_df : DataFrame with at least:
        - 'period_end' (or DateTimeIndex), 'real', 'exp',
        - weights: w_equities, w_gold, w_reits, w_bitcoin
        - kpis: kpi_eq, kpi_gold, kpi_reit, kpi_btc
    daily_returns : Optional daily returns (for benchmark aggregation)
        columns expected: ['Equities','Gold','REITs','Bitcoin']
    """
    # ---------- Prepare data (2023 only) ----------
    if "period_end" in log_df.columns:
        df = _ensure_dt_index(log_df, "period_end")
    else:
        df = log_df.copy()
        if not isinstance(df.index, pd.DatetimeIndex):
            raise ValueError("log_df needs a DatetimeIndex or a 'period_end' column.")
        df = df.sort_index()
        df.index = df.index.tz_localize(None)

    df = _year_slice(df, focus_year)

    # Strategy cumulative (monthly compounding)
    strat_cum = (1.0 + df["real"].fillna(0.0)).cumprod() - 1.0

    # Optional benchmark: NIFTY monthly from daily data, sliced to 2023
    bench_cum = None
    if isinstance(daily_returns, pd.DataFrame) and "Equities" in daily_returns.columns:
        dr = daily_returns.copy()
        dr.index = pd.to_datetime(dr.index).tz_localize(None)
        dr = _year_slice(dr, focus_year)
        if not dr.empty:
            eq_m = (1 + dr["Equities"]).resample("M").apply(lambda x: (1+x).prod()-1)
            eq_m.index = eq_m.index.to_period("M").to_timestamp("M")
            eq_m = eq_m.reindex(df.index, method="nearest")
            bench_cum = (1 + eq_m.fillna(0.0)).cumprod() - 1.0

    # Drawdowns
    running_max = (1 + strat_cum).cummax()
    drawdown = (1 + strat_cum) / running_max - 1.0

    # ---------- Summary metrics ----------
    total, vol, sharpe, maxdd = _summary_stats(df["real"])

    # ---------- Figure 1: Executive summary ----------
    fig0, ax0 = plt.subplots(figsize=(10, 5.6))
    ax0.axis("off")
    lines = [
        f"{title}",
        "",
        f"Total Return (2023): {100*total:.1f}%",
        f"Volatility (ann.): {100*vol:.1f}%",
        f"Sharpe (ann.): {sharpe:.2f}",
        f"Max Drawdown: {100*maxdd:.1f}%",
        "",
        "Key Context",
        "• Strong finish in Nov/Dec after a risk-off Q3.",
        "• Q3’23 saw rising US yields (10Y ↑), stronger USD, oil up → global risk-off.",
        "• Model stayed trend-anchored; KPIs toned down Equities risk as signals softened.",
    ]
    y = 0.92
    for i, ln in enumerate(lines):
        fz = 16 if i == 0 else (12 if ln and ln[0] != "•" else 10.5)
        ax0.text(0.02, y, ln, transform=ax0.transAxes, ha="left", va="top", fontsize=fz)
        y -= 0.06 if i == 0 else 0.08 if ln == "" else 0.055
    _add_q3_2023_band(ax0)  # (draws a subtle band behind text area; harmless if off-axes)
    fig0.tight_layout()
    fig0.savefig(f"{outfile_prefix}_00_summary.png", dpi=180)

    # ---------- Figure 2: Cumulative Returns (Strategy vs NIFTY) ----------
    fig1, ax1 = plt.subplots(figsize=(10, 5.6))
    ax1.plot(strat_cum.index, strat_cum.values, linewidth=2, label="Strategy")
    if isinstance(bench_cum, pd.Series):
        ax1.plot(bench_cum.index, bench_cum.values, linewidth=2, label="NIFTY")
    _add_q3_2023_band(ax1)
    ax1.set_title("Cumulative Return — 2023", fontsize=14)
    _pct_fmt(ax1)
    ax1.grid(True, alpha=0.25)
    ax1.legend(loc="upper left")
    fig1.tight_layout()
    fig1.savefig(f"{outfile_prefix}_01_cumulative.png", dpi=180)

    # ---------- Figure 3: Weights (stacked) ----------
    w_cols = ["w_equities", "w_gold", "w_reits", "w_bitcoin"]
    w_plot = df[w_cols].clip(lower=0, upper=1).fillna(0.0)
    fig2, ax2 = plt.subplots(figsize=(10, 5.6))
    ax2.stackplot(w_plot.index, w_plot.T.values, labels=["Equities","Gold","REITs","Bitcoin"], alpha=0.95)
    _add_q3_2023_band(ax2)
    ax2.set_ylim(0, 1)
    ax2.set_title("Portfolio Weights — 2023", fontsize=14)
    ax2.yaxis.set_major_formatter(lambda y, pos: f"{100*y:.0f}%")
    ax2.grid(True, alpha=0.25)
    ax2.legend(loc="upper left")
    fig2.tight_layout()
    fig2.savefig(f"{outfile_prefix}_02_weights.png", dpi=180)

    # ---------- Figure 4: Expected vs Realized (monthly) ----------
    fig3, ax3 = plt.subplots(figsize=(10, 5.6))
    # offset bars a few days left/right to create two clusters per month
    ax3.bar(df.index - pd.Timedelta(days=6), df["exp"].values, width=10, label="Expected", alpha=0.85)
    ax3.bar(df.index + pd.Timedelta(days=6), df["real"].values, width=10, label="Realized", alpha=0.85)
    _add_q3_2023_band(ax3)
    ax3.set_title("Expected vs Realized — 2023 (per month)", fontsize=14)
    _pct_fmt(ax3)
    ax3.grid(True, alpha=0.25)
    ax3.legend(loc="upper left")
    fig3.tight_layout()
    fig3.savefig(f"{outfile_prefix}_03_exp_vs_real.png", dpi=180)

    # ---------- Figure 5: Drawdown ----------
    fig4, ax4 = plt.subplots(figsize=(10, 5.0))
    ax4.fill_between(drawdown.index, drawdown.values, 0, step=None, alpha=0.35)
    _add_q3_2023_band(ax4)
    ax4.set_title("Drawdown — 2023", fontsize=14)
    _pct_fmt(ax4)
    ax4.grid(True, alpha=0.25)
    fig4.tight_layout()
    fig4.savefig(f"{outfile_prefix}_04_drawdown.png", dpi=180)

    # ---------- Figure 6: KPI dashboard (forward-looking signals) ----------
    kpi_cols = [("kpi_eq","Equities KPI"), ("kpi_gold","Gold KPI"), ("kpi_reit","REIT KPI"), ("kpi_btc","Bitcoin KPI")]
    have_any = any(c in df.columns for c,_ in kpi_cols)
    if have_any:
        fig5, ax5 = plt.subplots(figsize=(10, 5.6))
        for col, lbl in kpi_cols:
            if col in df.columns:
                ax5.plot(df.index, df[col].fillna(0.0).values, linewidth=2, label=lbl)
        _kpi_bands(ax5)
        _add_q3_2023_band(ax5)
        ax5.set_ylim(-1.05, 1.05)
        ax5.set_title("Forward-Looking KPI Signals — 2023", fontsize=14)
        ax5.grid(True, alpha=0.25)
        ax5.legend(loc="upper left")
        fig5.tight_layout()
        fig5.savefig(f"{outfile_prefix}_05_kpis.png", dpi=180)

    # ---------- Optional: 2-slide PPTX ----------
    if _PPTX_AVAILABLE:
        prs = Presentation()
        # Slide 1: Title + Summary + Cumulative + Weights
        s1 = prs.slides.add_slide(prs.slide_layouts[5])
        tx1 = s1.shapes.add_textbox(Inches(0.5), Inches(0.2), Inches(9.0), Inches(0.8))
        p = tx1.text_frame.paragraphs[0]
        p.text = title
        p.font.size = Pt(28); p.font.bold = True

        for i, img in enumerate([f"{outfile_prefix}_00_summary.png", f"{outfile_prefix}_01_cumulative.png", f"{outfile_prefix}_02_weights.png"]):
            if os.path.exists(img):
                top = 1.1 + i*2.7 if i > 0 else 1.1
                s1.shapes.add_picture(img, Inches(0.5), Inches(top), width=Inches(9.0))

        # Slide 2: Exp vs Real + Drawdown + KPIs
        s2 = prs.slides.add_slide(prs.slide_layouts[5])
        tx2 = s2.shapes.add_textbox(Inches(0.5), Inches(0.2), Inches(9.0), Inches(0.8))
        p2 = tx2.text_frame.paragraphs[0]
        p2.text = "Performance Diagnostics — 2023"
        p2.font.size = Pt(24); p2.font.bold = True

        y = 1.1
        for img in [f"{outfile_prefix}_03_exp_vs_real.png", f"{outfile_prefix}_04_drawdown.png", f"{outfile_prefix}_05_kpis.png"]:
            if os.path.exists(img):
                s2.shapes.add_picture(img, Inches(0.5), Inches(y), width=Inches(9.0))
                y += 2.8

        prs.save(f"{outfile_prefix}.pptx")
